from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.responses import FileResponse
from openai import OpenAI
from tavily import TavilyClient
from typing import List, Dict, Optional, Literal
from pydantic import BaseModel
import json
import PyPDF2
import io
import os
from dotenv import load_dotenv
from pptx import Presentation

load_dotenv()

# --- Pydantic Models ---
class Slide(BaseModel):
    title: str
    points: List[str]
    image_suggestion: Optional[str] = None

class PresentationRequest(BaseModel):
    outline: List[Slide]

# --- FastAPI App Initialization ---
app = FastAPI()

# --- Client Initializations ---
client = OpenAI(
    api_key=os.environ.get("OPENAI_API_KEY", "")
)
tavily_client = TavilyClient(api_key=os.environ.get("TAVILY_API_KEY", ""))


# --- Helper Functions ---
def extract_text_from_pdf(pdf_file: bytes) -> str:
    """Extracts text from a PDF file."""
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(pdf_file))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Lỗi đọc file PDF: {str(e)}")

def internet_search(
    query: str,
    max_results: int = 5,
    topic: Literal["general", "news", "finance"] = "general",
    include_raw_content: bool = False,
):
    """Run a web search for presentation research"""
    if not tavily_client.api_key:
        return {"error": "TAVILY_API_KEY not found in environment variables"}
    
    try:
        search_docs = tavily_client.search(
            query,
            max_results=max_results,
            include_raw_content=include_raw_content,
            topic=topic,
        )
        return search_docs
    except Exception as e:
        return {"error": f"Lỗi khi tìm kiếm: {str(e)}"}

def create_outline_prompt(topic: str, context: str, search_results: dict) -> str:
    """Creates a prompt for the LLM to generate a presentation outline."""
    return f"""
Bạn là một chuyên gia tạo nội dung thuyết trình. Nhiệm vụ của bạn là tạo ra một dàn ý chi tiết cho bài thuyết trình dựa trên chủ đề, tài liệu tham khảo và kết quả tìm kiếm trên internet.

**Chủ đề:** {topic}

**Tài liệu tham khảo (nếu có):**
{context}

**Kết quả tìm kiếm trên Internet:**
{json.dumps(search_results, indent=2)}

**Yêu cầu:**
1.  Phân tích chủ đề và các thông tin được cung cấp.
2.  Tạo một dàn ý (outline) cho bài thuyết trình gồm 5-10 slide.
3.  Mỗi slide trong dàn ý phải có:
    -   `title`: Tiêu đề của slide (ngắn gọn, súc tích).
    -   `points`: Một danh sách các gạch đầu dòng (bullet points) trình bày nội dung chính của slide.
    -   `image_suggestion`: Một gợi ý về hình ảnh minh họa cho slide.
4.  Trả về kết quả dưới dạng một JSON array, mỗi object là một slide.

**Ví dụ định dạng JSON:**
```json
[
    {{
        "title": "Tiêu đề Slide 1",
        "points": ["Nội dung 1.1", "Nội dung 1.2"],
        "image_suggestion": "Hình ảnh minh họa cho slide 1"
    }},
    {{
        "title": "Tiêu đề Slide 2",
        "points": ["Nội dung 2.1", "Nội dung 2.2"],
        "image_suggestion": "Biểu đồ minh họa cho slide 2"
    }}
]
```
Hãy đảm bảo dàn ý logic, mạch lạc và bao quát được chủ đề.
"""

def create_slide_content_prompt(slide_title: str, outline_points: List[str], search_results: dict) -> str:
    """Creates a prompt for the LLM to generate detailed slide content."""
    return f"""
Bạn là một chuyên gia tạo nội dung thuyết trình. Nhiệm vụ của bạn là viết nội dung chi tiết cho một slide dựa trên tiêu đề, các điểm chính trong dàn ý và kết quả tìm kiếm.

**Tiêu đề Slide:** {slide_title}

**Các điểm chính từ dàn ý:**
- {"\n- ".join(outline_points)}

**Kết quả tìm kiếm trên Internet để tham khảo:**
{json.dumps(search_results, indent=2)}

**Yêu cầu:**
1.  Dựa vào thông tin trên, hãy viết lại nội dung cho slide một cách chi tiết và hấp dẫn.
2.  Giữ lại tiêu đề gốc.
3.  Mở rộng các điểm chính thành những câu văn hoàn chỉnh, cung cấp thêm thông tin, ví dụ, hoặc số liệu nếu có.
4.  Trả về kết quả dưới dạng một JSON object với hai key: `title` (giữ nguyên) và `points` (danh sách nội dung chi tiết mới).

**Ví dụ định dạng JSON:**
```json
{{
    "title": "{slide_title}",
    "points": [
        "Nội dung chi tiết cho điểm 1, được viết lại đầy đủ hơn.",
        "Nội dung chi tiết cho điểm 2, có thể bao gồm số liệu từ kết quả tìm kiếm.",
        "Một điểm mới được bổ sung nếu cần thiết để làm rõ ý."
    ]
}}
```
"""

# --- API Endpoints ---
@app.post("/generate-outline")
async def generate_outline(topic: str = Form(...), file: Optional[UploadFile] = File(None)):
    """
    Generates a presentation outline from a topic, optional PDF, and internet search.
    """
    # 1. Extract context from PDF if provided
    pdf_context = ""
    if file:
        if not file.filename.lower().endswith('.pdf'):
            raise HTTPException(status_code=400, detail="Chỉ chấp nhận file PDF")
        try:
            pdf_content = await file.read()
            pdf_context = extract_text_from_pdf(pdf_content)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Lỗi xử lý file: {str(e)}")

    # 2. Perform internet search
    search_results = internet_search(query=f"Outline for presentation on {topic}")
    if "error" in search_results:
        raise HTTPException(status_code=500, detail=search_results["error"])

    # 3. Generate outline using LLM
    prompt = create_outline_prompt(topic, pdf_context, search_results)
    try:
        response = client.chat.completions.create(
            model="gpt-4.1",
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"}
        )
        outline_json = json.loads(response.choices[0].message.content)
        return outline_json
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi khi gọi LLM để tạo dàn ý: {str(e)}")


@app.post("/generate-presentation")
async def generate_presentation(request: PresentationRequest):
    """
    Generates a PPTX file from a given outline, enriching each slide with search results.
    """
    prs = Presentation()
    
    for slide_data in request.outline:
        # 1. Search for detailed content for the current slide
        query = f"Detailed information for a presentation slide titled '{slide_data.title}' covering points: {', '.join(slide_data.points)}"
        search_results = internet_search(query, max_results=3)
        if "error" in search_results:
            # Continue without search results if search fails
            search_results = {}

        # 2. Generate detailed slide content using LLM
        prompt = create_slide_content_prompt(slide_data.title, slide_data.points, search_results)
        try:
            response = client.chat.completions.create(
                model="gpt-4.1",
                messages=[{"role": "user", "content": prompt}],
                response_format={"type": "json_object"}
            )
            slide_content = json.loads(response.choices[0].message.content)
            final_title = slide_content.get("title", slide_data.title)
            final_points = slide_content.get("points", slide_data.points)
        except Exception:
            # If LLM fails, use original content
            final_title = slide_data.title
            final_points = slide_data.points

        # 3. Add the enriched content to the presentation slide
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = final_title
        
        tf = content_shape.text_frame
        tf.clear() 
        
        for point in final_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 1

    output_path = "generated_presentation.pptx"
    prs.save(output_path)
    
    return FileResponse(
        output_path, 
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', 
        filename='presentation.pptx'
    )

# --- Main Execution ---
if __name__ == "__main__":
    import uvicorn
    uvicorn.run("agent_pptx_generator:app", host="0.0.0.0", port=8002, reload=True)
